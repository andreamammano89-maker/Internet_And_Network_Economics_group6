from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ── Color palette ──────────────────────────────────────────────────
BLUE       = RGBColor(0x18, 0x5F, 0xA5)
DARK_NAVY  = RGBColor(0x0C, 0x44, 0x7C)
LIGHT_BLUE = RGBColor(0xE6, 0xF1, 0xFB)
RED_IT     = RGBColor(0x99, 0x3C, 0x1D)
LIGHT_RED  = RGBColor(0xFA, 0xEC, 0xE7)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BODY_TEXT  = RGBColor(0x2C, 0x2C, 0x2A)
MUTED      = RGBColor(0x88, 0x87, 0x80)
GRAY_FILL  = RGBColor(0xF1, 0xEF, 0xE8)
GRAY_BRD   = RGBColor(0x88, 0x87, 0x80)
DARK_BLUE_FILL = RGBColor(0x18, 0x5F, 0xA5)

FONT = "Calibri"
MARGIN = Inches(0.4)
FOOTER_TEXT = "Group 6 | Internet and Network Economics | LUISS"

# ── Presentation setup ─────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]   # completely blank


# ══════════════════════════════════════════════════════════════════════
#  Helper functions
# ══════════════════════════════════════════════════════════════════════

def rgb_hex(r):
    return "{:02X}{:02X}{:02X}".format(r[0], r[1], r[2])

def set_bg(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill=None, line=None, line_width_pt=1.5, radius=None):
    """Add a rectangle shape and return it."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        x, y, w, h
    )
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is not None:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    if radius is not None:
        # set rounded corners via XML
        sp = shape._element
        prstGeom = sp.find(qn('p:spPr') + '/' + qn('a:prstGeom'))
        if prstGeom is None:
            spPr = sp.find(qn('p:spPr'))
            prstGeom = etree.SubElement(spPr, qn('a:prstGeom'))
            prstGeom.set('prst', 'roundRect')
        else:
            prstGeom.set('prst', 'roundRect')
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        for gd in avLst.findall(qn('a:gd')):
            avLst.remove(gd)
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', f'val {radius}')
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=12, bold=False, italic=False,
             color=BODY_TEXT, align=PP_ALIGN.LEFT,
             wrap=True, font_name=FONT):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_textbox_rich(slide, lines, x, y, w, h,
                     wrap=True, bg_fill=None, border_color=None,
                     border_width=1.5, padding=Inches(0.12)):
    """
    lines: list of dicts with keys text, size, bold, italic, color, align
    """
    if bg_fill or border_color:
        shape = add_rect(slide, x, y, w, h,
                         fill=bg_fill,
                         line=border_color,
                         line_width_pt=border_width)
        txb = shape
        tf = shape.text_frame
    else:
        txb = slide.shapes.add_textbox(x, y, w, h)
        tf = txb.text_frame
    tf.word_wrap = wrap
    # clear default empty paragraph
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = line.get('align', PP_ALIGN.LEFT)
        if line.get('space_before'):
            p.space_before = Pt(line['space_before'])
        run = p.add_run()
        run.text = line['text']
        run.font.name = line.get('font', FONT)
        run.font.size = Pt(line.get('size', 12))
        run.font.bold = line.get('bold', False)
        run.font.italic = line.get('italic', False)
        run.font.color.rgb = line.get('color', BODY_TEXT)
    return txb

def add_footer(slide):
    add_text(slide, FOOTER_TEXT,
             MARGIN, Inches(7.15),
             Inches(12.5), Inches(0.3),
             font_size=8, color=MUTED, align=PP_ALIGN.CENTER)

def accent_bar(slide):
    """Add a blue left accent bar."""
    add_rect(slide, Inches(0), Inches(0),
             Inches(0.1), prs.slide_height,
             fill=BLUE)

def slide_title(slide, title_text):
    add_text(slide, title_text.upper(),
             MARGIN + Inches(0.15), Inches(0.3),
             Inches(12.4), Inches(0.7),
             font_size=28, bold=True, color=DARK_NAVY)

def placeholder_box(slide, x, y, w, h, label):
    add_rect(slide, x, y, w, h,
             fill=LIGHT_BLUE, line=BLUE, line_width_pt=1.5)
    add_text(slide, label,
             x + Inches(0.15), y + Inches(0.15),
             w - Inches(0.3), h - Inches(0.3),
             font_size=11, color=BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
#  SLIDE 1 – Cover
# ══════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(blank_layout)
set_bg(slide1, BLUE)

# Decorative white rounded rectangle top-right (3×3 in, partially off)
add_rect(slide1,
         Inches(11.2), Inches(-0.8),
         Inches(3.0), Inches(3.0),
         fill=WHITE, radius=20000)

# Title
add_textbox_rich(slide1, [
    {'text': 'INTANGIBLE CAPITAL AND\nLABOUR PRODUCTIVITY\nDIVERGENCE',
     'size': 40, 'bold': True, 'color': WHITE, 'align': PP_ALIGN.LEFT},
], Inches(0.5), Inches(1.6), Inches(9.5), Inches(3.2))

# Subtitle
add_text(slide1, "Denmark vs Italy  ·  ICT Sector  ·  1996–2019",
         Inches(0.5), Inches(4.95), Inches(9.0), Inches(0.5),
         font_size=18, color=WHITE, align=PP_ALIGN.LEFT)

# Bottom credit
add_text(slide1, "Group 6  |  Internet and Network Economics  |  LUISS 2025",
         Inches(0.5), Inches(6.8), Inches(9.0), Inches(0.4),
         font_size=11, color=WHITE, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════
#  SLIDE 2 – The Puzzle
# ══════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(blank_layout)
accent_bar(slide2)
slide_title(slide2, "The productivity puzzle")
add_footer(slide2)

# Left text box
add_textbox_rich(slide2, [
    {'text': "Europe's productivity stagnation",
     'size': 13, 'bold': True, 'color': DARK_NAVY},
    {'text': "\nHigh-income economies have experienced underwhelming productivity growth since the 1990s. Digital technologies were expected to reverse this trend — but their impact is highly uneven across countries.",
     'size': 11, 'bold': False, 'color': BODY_TEXT},
], Inches(0.6), Inches(1.2), Inches(5.6), Inches(4.0))

# Right research question box
add_textbox_rich(slide2, [
    {'text': "Research Question",
     'size': 12, 'bold': True, 'color': DARK_NAVY},
    {'text': "\n\"To what extent do differences in the contribution of intangible capital services explain divergence in labour productivity growth between Denmark and Italy over the period 1995–2019?\"",
     'size': 11, 'bold': False, 'italic': True, 'color': BODY_TEXT},
], Inches(6.6), Inches(1.2), Inches(6.3), Inches(4.0),
   bg_fill=LIGHT_BLUE, border_color=BLUE)


# ══════════════════════════════════════════════════════════════════════
#  SLIDE 3 – Theoretical Framework
# ══════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(blank_layout)
accent_bar(slide3)
slide_title(slide3, "Growth accounting framework")
add_footer(slide3)

boxes = [
    (LIGHT_BLUE, BLUE, "Δ Intangible capital", "LP1ConIntang"),
    (GRAY_FILL,  GRAY_BRD, "Δ Tangible ICT", "LP1ConTangICT"),
    (GRAY_FILL,  GRAY_BRD, "Δ TFP", "LP1ConTFP"),
]
box_w = Inches(2.4)
box_h = Inches(1.3)
box_y = Inches(1.9)
gap   = Inches(0.3)
start_x = Inches(0.6)

for i, (fill, border, label, sub) in enumerate(boxes):
    bx = start_x + i * (box_w + gap + Inches(0.3))
    add_rect(slide3, bx, box_y, box_w, box_h, fill=fill, line=border)
    add_textbox_rich(slide3, [
        {'text': label, 'size': 12, 'bold': True, 'color': DARK_NAVY, 'align': PP_ALIGN.CENTER},
        {'text': sub, 'size': 9, 'bold': False, 'color': MUTED, 'align': PP_ALIGN.CENTER},
    ], bx + Inches(0.1), box_y + Inches(0.1),
       box_w - Inches(0.2), box_h - Inches(0.2))

    # plus sign between boxes
    if i < 2:
        add_text(slide3, "+",
                 bx + box_w + Inches(0.05),
                 box_y + Inches(0.35),
                 Inches(0.2), Inches(0.6),
                 font_size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# equals sign
eq_x = start_x + 3 * (box_w + gap + Inches(0.3)) - Inches(0.3)
add_text(slide3, "=",
         eq_x, box_y + Inches(0.35),
         Inches(0.3), Inches(0.6),
         font_size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# Result box
res_x = eq_x + Inches(0.4)
add_rect(slide3, res_x, box_y, box_w, box_h, fill=BLUE)
add_textbox_rich(slide3, [
    {'text': "LP Growth", 'size': 13, 'bold': True, 'color': WHITE, 'align': PP_ALIGN.CENTER},
    {'text': "LP1_G", 'size': 9, 'bold': False, 'color': WHITE, 'align': PP_ALIGN.CENTER},
], res_x + Inches(0.1), box_y + Inches(0.15),
   box_w - Inches(0.2), box_h - Inches(0.3))

# Description
add_textbox_rich(slide3, [
    {'text': "Under competitive markets and constant returns to scale, labour productivity growth equals the weighted sum of capital deepening contributions plus TFP growth (Solow, 1957; Jorgenson & Griliches, 1967).",
     'size': 11, 'color': BODY_TEXT},
    {'text': "\nSource: EUKLEMS & INTANProd (2024). Bontadini et al. (2023).",
     'size': 9, 'italic': True, 'color': MUTED},
], Inches(0.6), Inches(3.5), Inches(12.2), Inches(1.6))


# ══════════════════════════════════════════════════════════════════════
#  SLIDE 4 – Data & Methodology
# ══════════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(blank_layout)
accent_bar(slide4)
slide_title(slide4, "Data & methodology")
add_footer(slide4)

cols = [
    ("Dataset",
     "EUKLEMS & INTANProd 2024 release.\nGrowth Accounts Extended + Intangibles Analytical."),
    ("Variables",
     "LP1_G · LP1ConIntang · LP1ConTFP · LP1ConTangICT · I_Intang · I_OrgCap · I_Soft_DB"),
    ("Scope",
     "Countries: Denmark, Italy\nSector: J (ICT) primary; TOT_IND benchmark\nPeriod: 1996–2019 (COVID excluded)"),
]
col_w = Inches(3.9)
col_h = Inches(4.2)
col_y = Inches(1.3)
col_gap = Inches(0.35)
col_start = Inches(0.6)

for i, (hdr, body) in enumerate(cols):
    cx = col_start + i * (col_w + col_gap)
    # header
    add_rect(slide4, cx, col_y, col_w, Inches(0.55), fill=BLUE)
    add_text(slide4, hdr,
             cx + Inches(0.1), col_y + Inches(0.05),
             col_w - Inches(0.2), Inches(0.45),
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # body
    add_textbox_rich(slide4, [
        {'text': body, 'size': 11, 'color': BODY_TEXT},
    ], cx, col_y + Inches(0.55), col_w, col_h - Inches(0.55),
       bg_fill=LIGHT_BLUE, border_color=BLUE)


# ══════════════════════════════════════════════════════════════════════
#  SLIDE 5 – The Divergence
# ══════════════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(blank_layout)
accent_bar(slide5)
slide_title(slide5, "Labour productivity growth: Denmark outperforms Italy")
add_footer(slide5)

placeholder_box(slide5,
                Inches(0.6), Inches(1.2),
                Inches(8.5), Inches(5.3),
                "[INSERT Chart 1: LP1_G line chart, sector J, 1996–2019]")

# Key finding box
add_textbox_rich(slide5, [
    {'text': "Key Finding", 'size': 12, 'bold': True, 'color': WHITE, 'align': PP_ALIGN.CENTER},
    {'text': "\nAverage gap:\n+3.35 p.p. per year\nin favour of Denmark\nacross the full\n1996–2019 period",
     'size': 11, 'color': WHITE, 'align': PP_ALIGN.CENTER},
], Inches(9.4), Inches(1.2), Inches(3.5), Inches(5.3),
   bg_fill=BLUE, border_color=DARK_NAVY)


# ══════════════════════════════════════════════════════════════════════
#  SLIDE 6 – Gap Decomposition
# ══════════════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(blank_layout)
accent_bar(slide6)
slide_title(slide6, "Decomposing the gap: TFP dominates, intangibles play a structural role")
add_footer(slide6)

placeholder_box(slide6,
                Inches(0.6), Inches(1.2),
                Inches(12.3), Inches(4.1),
                "[INSERT Chart 2: Stacked bar decomposition, sector J, annual]")

stats = [
    (BLUE,       WHITE,      "~3.2 pp",  "TFP contribution (avg)"),
    (LIGHT_BLUE, DARK_NAVY,  "~0.5 pp",  "Intangible contribution (avg)"),
    (GRAY_FILL,  BODY_TEXT,  "~−0.2 pp", "ICT tangible contribution (avg)"),
]
stat_w = Inches(3.9)
stat_y = Inches(5.6)
for i, (fill, tc, big, small) in enumerate(stats):
    sx = Inches(0.6) + i * (stat_w + Inches(0.25))
    add_rect(slide6, sx, stat_y, stat_w, Inches(1.1), fill=fill)
    add_textbox_rich(slide6, [
        {'text': big,   'size': 15, 'bold': True,  'color': tc, 'align': PP_ALIGN.CENTER},
        {'text': small, 'size': 10, 'bold': False, 'color': tc, 'align': PP_ALIGN.CENTER},
    ], sx + Inches(0.1), stat_y + Inches(0.1),
       stat_w - Inches(0.2), Inches(0.9))


# ══════════════════════════════════════════════════════════════════════
#  SLIDE 7 – Intangible Channel
# ══════════════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(blank_layout)
accent_bar(slide7)
slide_title(slide7, "The intangible channel: positive before 2012, negative after for Italy")
add_footer(slide7)

placeholder_box(slide7,
                Inches(0.6), Inches(1.2),
                Inches(5.9), Inches(3.8),
                "[INSERT Chart 6: Sub-period decomposition]")
placeholder_box(slide7,
                Inches(6.9), Inches(1.2),
                Inches(6.0), Inches(3.8),
                "[INSERT Chart 3: LP1ConIntang over time, DK vs IT]")

add_textbox_rich(slide7, [
    {'text': "Key finding: ", 'size': 11, 'bold': True, 'color': RED_IT},
    {'text': "after 2011, Italy's intangible capital contribution turns negative while Denmark's remains positive — a structural divergence within the divergence.",
     'size': 11, 'bold': False, 'color': BODY_TEXT},
], Inches(0.6), Inches(5.25), Inches(12.3), Inches(1.1),
   bg_fill=LIGHT_RED, border_color=RED_IT)


# ══════════════════════════════════════════════════════════════════════
#  SLIDE 8 – Investment and Efficiency
# ══════════════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(blank_layout)
accent_bar(slide8)
slide_title(slide8, "Denmark invests more and converts it more efficiently")
add_footer(slide8)

placeholder_box(slide8,
                Inches(0.6), Inches(1.2),
                Inches(5.9), Inches(3.9),
                "[INSERT Chart 4: I_Intang levels + OrgCap composition]")
placeholder_box(slide8,
                Inches(6.9), Inches(1.2),
                Inches(6.0), Inches(3.9),
                "[INSERT Chart 7: Conversion efficiency LP1ConIntang / I_Intang]")

add_text(slide8,
         "Conversion efficiency = LP1ConIntang / I_Intang × 1000.  "
         "A 3-year centred rolling average is applied to smooth annual volatility.",
         Inches(0.6), Inches(5.3), Inches(12.3), Inches(0.7),
         font_size=10, italic=True, color=MUTED)


# ══════════════════════════════════════════════════════════════════════
#  SLIDE 9 – Benchmark
# ══════════════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(blank_layout)
accent_bar(slide9)
slide_title(slide9, "The ICT sector gap exceeds the total economy gap")
add_footer(slide9)

placeholder_box(slide9,
                Inches(0.6), Inches(1.2),
                Inches(5.9), Inches(3.7),
                "[INSERT Chart 5 left: LP1_G TOT_IND, DK vs IT]")
placeholder_box(slide9,
                Inches(6.9), Inches(1.2),
                Inches(6.0), Inches(3.7),
                "[INSERT Chart 5 right: Gap decomposition TOT_IND]")

add_textbox_rich(slide9, [
    {'text': "The divergence is not a general macroeconomic phenomenon — it is concentrated in the ICT sector, where intangible assets matter most.",
     'size': 11, 'bold': False, 'color': WHITE, 'align': PP_ALIGN.CENTER},
], Inches(0.6), Inches(5.2), Inches(12.3), Inches(1.05),
   bg_fill=BLUE)


# ══════════════════════════════════════════════════════════════════════
#  SLIDE 10 – Institutional Layer
# ══════════════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(blank_layout)
accent_bar(slide10)
slide_title(slide10, "Institutional context: why Denmark converts better")
add_footer(slide10)

# --- Table ---
headers = ["Indicator", "Denmark", "Italy", "Source"]
rows = [
    ["Number of enterprises, Sector J (2019)", "18,961", "104,879", "Eurostat SBS"],
    ["Avg. persons employed per enterprise, Sector J (2019)", "6.5", "5.6", "Eurostat SBS"],
    ["EPL Index — regular contracts (2019)", "1.53", "2.56", "OECD EPL Database"],
]

col_widths = [Inches(4.8), Inches(1.8), Inches(1.8), Inches(3.2)]
row_h = Inches(0.48)
tbl_x = Inches(0.6)
tbl_y = Inches(1.15)
tbl_w = sum(col_widths)

# header row
cx = tbl_x
for j, hdr in enumerate(headers):
    add_rect(slide10, cx, tbl_y, col_widths[j], row_h, fill=BLUE)
    add_text(slide10, hdr, cx + Inches(0.08), tbl_y + Inches(0.06),
             col_widths[j] - Inches(0.16), row_h - Inches(0.1),
             font_size=11, bold=True, color=WHITE)
    cx += col_widths[j]

# data rows
for i, row in enumerate(rows):
    ry = tbl_y + (i + 1) * row_h
    fill = WHITE if i % 2 == 0 else LIGHT_BLUE
    cx = tbl_x
    for j, cell in enumerate(row):
        add_rect(slide10, cx, ry, col_widths[j], row_h, fill=fill, line=BLUE, line_width_pt=0.5)
        is_bold = (j == 1)   # bold Denmark values
        add_text(slide10, cell,
                 cx + Inches(0.08), ry + Inches(0.06),
                 col_widths[j] - Inches(0.16), row_h - Inches(0.1),
                 font_size=10, bold=is_bold, color=BODY_TEXT)
        cx += col_widths[j]

# Two text boxes below table
tbl_bottom = tbl_y + (len(rows) + 1) * row_h + Inches(0.15)
half_w = Inches(6.1)

add_textbox_rich(slide10, [
    {'text': "Flexicurity model:", 'size': 11, 'bold': True, 'color': DARK_NAVY},
    {'text': " Denmark's labour market flexibility lowers the cost of organisational restructuring — a prerequisite for effective intangible capital investment.",
     'size': 11, 'bold': False, 'color': BODY_TEXT},
], Inches(0.6), tbl_bottom, half_w, Inches(1.15),
   border_color=BLUE)

add_textbox_rich(slide10, [
    {'text': "Fragmented firm structure:", 'size': 11, 'bold': True, 'color': RED_IT},
    {'text': " Italy's ~105,000 ICT enterprises face structural barriers to investing in organisational capital at scale.",
     'size': 11, 'bold': False, 'color': BODY_TEXT},
], Inches(7.05), tbl_bottom, half_w, Inches(1.15),
   border_color=RED_IT)

# Causality note
add_text(slide10,
         "Causality note: institutional factors are interpreted as consistent with the data patterns. "
         "No formal causal identification is established.",
         Inches(0.6), Inches(6.75), Inches(12.3), Inches(0.35),
         font_size=8, italic=True, color=MUTED)


# ══════════════════════════════════════════════════════════════════════
#  SLIDE 11 – Answer to Research Question
# ══════════════════════════════════════════════════════════════════════
slide11 = prs.slides.add_slide(blank_layout)
set_bg(slide11, BLUE)

add_text(slide11, "ANSWER TO THE RESEARCH QUESTION",
         Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.65),
         font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rq_text = ('"To what extent do differences in the contribution of intangible capital services '
           'explain divergence in labour productivity growth between Denmark and Italy '
           'over the period 1995–2019?"')
add_text(slide11, rq_text,
         Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.85),
         font_size=11, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

answer = (
    "Intangible capital explains approximately 14% of the annual LP growth gap in the ICT sector "
    "(≈0.46 of 3.35 p.p.). TFP remains the dominant driver across all sub-periods. However, the "
    "intangible channel is not static: it is positive and growing before 2012, then turns negative "
    "for Italy while remaining positive for Denmark — a structural divergence that compounds the TFP "
    "gap. This pattern is consistent with institutional differences in labour market flexibility and "
    "firm size that limit Italy's capacity to convert intangible investment into productivity gains."
)
add_textbox_rich(slide11, [
    {'text': answer, 'size': 12, 'color': DARK_NAVY, 'align': PP_ALIGN.LEFT},
], Inches(0.7), Inches(2.1), Inches(11.9), Inches(4.5),
   bg_fill=WHITE, border_color=BLUE, border_width=2.0)


# ══════════════════════════════════════════════════════════════════════
#  SLIDE 12 – Conclusions
# ══════════════════════════════════════════════════════════════════════
slide12 = prs.slides.add_slide(blank_layout)
accent_bar(slide12)
slide_title(slide12, "Conclusions & limitations")
add_footer(slide12)

left_w  = Inches(7.5)
right_w = Inches(4.7)
col_y12 = Inches(1.15)

# LEFT – Conclusions
conc_lines = [
    {'text': "Conclusions", 'size': 13, 'bold': True, 'color': DARK_NAVY},
    {'text': "\n· Denmark's ICT sector LP advantage over Italy is persistent, averaging +3.35 p.p. annually over 1996–2019.",
     'size': 11, 'color': BODY_TEXT},
    {'text': "\n· Intangibles explain ~14% of the gap directly, but their declining contribution in Italy post-2011 signals a structural deterioration.",
     'size': 11, 'color': BODY_TEXT},
    {'text': "\n· The conversion efficiency gap — not just investment volume — is the most original finding: institutional flexibility enables Denmark to extract more LP growth per euro of intangible investment.",
     'size': 11, 'color': BODY_TEXT},
]
add_textbox_rich(slide12, conc_lines,
                 Inches(0.6), col_y12, left_w, Inches(4.8))

# RIGHT – Limitations
lim_lines = [
    {'text': "Limitations", 'size': 13, 'bold': True, 'color': MUTED},
    {'text': "\n· Aggregate sector-level data: no firm-level identification possible.",
     'size': 10, 'color': MUTED},
    {'text': "\n· No formal causal identification: institutional links are interpretive.",
     'size': 10, 'color': MUTED},
    {'text': "\n· Other factors (R&D policy, digital infrastructure, education) not captured.",
     'size': 10, 'color': MUTED},
]
add_textbox_rich(slide12, lim_lines,
                 Inches(8.4), col_y12, right_w, Inches(4.8))

# Source note
add_text(slide12,
         "Data: EUKLEMS & INTANProd (2024). Eurostat SBS. OECD EPL Database.",
         Inches(0.6), Inches(6.45), Inches(12.0), Inches(0.35),
         font_size=9, italic=True, color=MUTED)

# Thank you
add_text(slide12, "Thank you",
         Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.45),
         font_size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
#  Save
# ══════════════════════════════════════════════════════════════════════
output_path = "DK_IT_Presentation.pptx"
prs.save(output_path)
print(f"Done! Presentation saved as {output_path}")
